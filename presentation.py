import players as pl

from pptx import Presentation
from pptx.util import Emu, Inches

from pathlib import Path
from random import shuffle, choice, sample
import re
import time

from itertools import zip_longest, chain

def sanitize_filename(name: str) -> str:
    name = name.strip()
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', '', name)
    name = name.strip('. ')
    return name or f"I tried to break the game {time.time_ns()}"



presentations_path = Path("./presentations")

def generate_presentation(player: pl.Player, title_text: str, subtitle_text: str, slide_paths: list):
    prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = title_text
    subtitle.text = f"{subtitle_text}\n\n{player.first_name} {player.last_name}"
    
    blank_slide_layout = prs.slide_layouts[6]
    for slide_path in slide_paths:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        pic = slide.shapes.add_picture(str(slide_path), 0, 0, width=prs.slide_width)
        if pic.height > prs.slide_height:
            pic.width, pic.height = Emu(pic.width * prs.slide_height / pic.height), prs.slide_height
            
        pic.left = Emu((prs.slide_width - pic.width) / 2)
        pic.top = Emu((prs.slide_height - pic.height) / 2)
        
    end_slide_layout = title_slide_layout
    end_slide = prs.slides.add_slide(end_slide_layout)
    
    thanks = end_slide.shapes.title
    name_surname = end_slide.placeholders[1]
    thanks.text = "Thank you for your attention."
    name_surname.text = f"{player.first_name} {(player.last_name or '')}"
    
    prs_path = presentations_path / f"{sanitize_filename(player.first_name)}.pptx"
    prs.save(prs_path)
    


def validate_players(players: pl.Players):
    total_slides = sum(player.count_slides for player in players.values())
    if total_slides < len(players):
        raise ValueError(f"fewer slides ({total_slides}) than players ({len(players)})")
    
    total_titles = sum(player.count_titles for player in players.values())
    if total_titles < len(players):
        raise ValueError(f"fewer titles ({total_titles}) than players ({len(players)})")
    
    total_subtitles = sum(player.count_subtitles for player in players.values())
    if total_subtitles < len(players):
        raise ValueError(f"fewer subtitles ({total_subtitles}) than players ({len(players)})")



def generate_all(players: pl.Players):
    validate_players(players)
    
    all_slide_paths = []
    all_titles = []
    all_subtitles = []
    for id, player in players.items():
        shuffle(player.slide_paths)
        shuffle(player.titles)
        shuffle(player.subtitles)
        
        # [ [(id0, slide1), (id0, slide0)], [(id1, slide2), ...], ... ]
        all_slide_paths.append( [ (id, i) for i in player.slide_paths[::-1] ] )
        all_titles.append( [ (id, i) for i in player.titles[::-1] ] )
        all_subtitles.append( [ (id, i) for i in player.subtitles[::-1] ] )
    
    for lst in (all_slide_paths, all_titles, all_subtitles):
        shuffle(lst)
    
    # roll em out
    # [(id0, slide1), (id1, slide2), ..., (id0, slide0), ...]
    all_slide_paths, all_titles, all_subtitles = (
        [pair for pair
            in (
                chain.from_iterable(zip_longest(*lst))
            )
            if pair is not None]
        for lst in (all_slide_paths, all_titles, all_subtitles)
    )
    
    titles = {}
    subtitles = {}
    for id, player in players.items():
        title = next(i for i in all_titles if i[0] != id)   # get from others
        all_titles.remove(title)
        titles[id] = title[1]
        
        subtitle = next(i for i in all_subtitles if i[0] != id)
        all_subtitles.remove(subtitle)
        subtitles[id] = subtitle[1]
    
    slide_paths = {id: [] for id in players.keys()}
    while len(all_slide_paths) >= len(players):
        for id, player in players.items():
            slide_path = next(i for i in all_slide_paths if i[0] != id)
            all_slide_paths.remove(slide_path)
            slide_paths[id].append(slide_path[1])
            
    for id, player in players.items():
        print(id)
        generate_presentation(player, titles[id], subtitles[id], slide_paths[id])
    


if __name__ == "__main__":
    generate_all(pl.players)